from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme - Modern IT theme with gradients
DARK_BLUE = RGBColor(25, 75, 150)
LIGHT_BLUE = RGBColor(70, 150, 220)
ACCENT_GREEN = RGBColor(76, 175, 80)
ACCENT_ORANGE = RGBColor(255, 140, 0)
DARK_GRAY = RGBColor(60, 60, 60)
LIGHT_GRAY = RGBColor(240, 240, 240)
WHITE = RGBColor(255, 255, 255)
TECH_NAVY = RGBColor(12, 44, 84)
TECH_CYAN = RGBColor(0, 188, 212)

def add_gradient_background(slide, color1, color2):
    """Add gradient background to a slide"""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 45.0
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2

def add_tech_pattern(slide, width, height):
    """Add subtle tech pattern to background"""
    # Add diagonal lines pattern
    for i in range(0, int(width * 914400), int(1 * 914400)):
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(i/914400), 
            Inches(-0.5), 
            Inches(0.01), 
            Inches(height + 1)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(255, 255, 255)
        line.fill.transparency = 0.95
        line.line.color.rgb = RGBColor(255, 255, 255)
        line.line.transparency = 0.95

def add_title_slide(prs, title, subtitle):
    """Add a title slide with tech background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_gradient_background(slide, TECH_NAVY, DARK_BLUE)
    
    # Add decorative circles
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7), Inches(-1), Inches(3), Inches(3))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = TECH_CYAN
    circle1.fill.transparency = 0.85
    circle1.line.color.rgb = TECH_CYAN
    
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0), Inches(5), Inches(2.5), Inches(2.5))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = ACCENT_GREEN
    circle2.fill.transparency = 0.85
    circle2.line.color.rgb = ACCENT_GREEN
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(66)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = subtitle
    subtitle_p.font.size = Pt(32)
    subtitle_p.font.color.rgb = TECH_CYAN
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # Add decorative line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(2.3), Inches(4), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_GREEN
    line.line.color.rgb = ACCENT_GREEN

def add_content_slide(prs, title, content_list, with_bullet=True):
    """Add a content slide with IT theme"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_gradient_background(slide, WHITE, LIGHT_GRAY)
    
    # Add header bar with gradient
    header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    header_fill = header_shape.fill
    header_fill.gradient()
    header_fill.gradient_stops[0].color.rgb = DARK_BLUE
    header_fill.gradient_stops[1].color.rgb = TECH_NAVY
    header_shape.line.color.rgb = ACCENT_GREEN
    header_shape.line.width = Pt(2)
    
    # Add accent bar on left
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), Inches(7.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT_GREEN
    accent_bar.line.color.rgb = ACCENT_GREEN
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5.2))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        
        if with_bullet:
            p.level = 0
            p.font.bold = False
    
    return slide

def add_two_column_slide(prs, title, left_content, right_content):
    """Add a two-column content slide with IT theme"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, WHITE, LIGHT_GRAY)
    
    # Add header bar with gradient
    header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
    header_fill = header_shape.fill
    header_fill.gradient()
    header_fill.gradient_stops[0].color.rgb = DARK_BLUE
    header_fill.gradient_stops[1].color.rgb = TECH_NAVY
    header_shape.line.color.rgb = ACCENT_GREEN
    header_shape.line.width = Pt(2)
    
    # Add accent bar
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), Inches(7.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT_GREEN
    accent_bar.line.color.rgb = ACCENT_GREEN
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(5.7))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(left_content):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.3), Inches(4.2), Inches(5.7))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(right_content):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    return slide

# ========== SLIDE 1: Title Slide ==========
add_title_slide(prs, "Playwright Framework", "Automated End-to-End Testing Solution")

# ========== SLIDE 2: Overview ==========
overview_content = [
    "🎯 Comprehensive end-to-end testing framework",
    "🎬 Built for movie ticket booking platform",
    "🏗️ Follows enterprise best practices",
    "⚡ Scalable, maintainable, and efficient",
    "👥 Supports cross-browser and parallel testing"
]
add_content_slide(prs, "Framework Overview", overview_content, with_bullet=True)

# ========== SLIDE 3: Tech Stack ==========
left_tech = [
    "✓ Playwright 1.58.2",
    "✓ TypeScript",
    "✓ Node.js",
    "✓ Allure Reporting"
]
right_tech = [
    "✓ Faker.js (Test Data Generation)",
    "✓ ESM Modules",
    "✓ Dotenv (Environment Config)",
    "✓ Cross-Browser Support"
]
add_two_column_slide(prs, "Technology Stack", left_tech, right_tech)

# ========== SLIDE 4: Supported Browsers ==========
browsers_content = [
    "🌐 Chromium - Fastest, production-ready",
    "🔥 Firefox - Security & standards compliance",
    "📱 WebKit - Mobile browser simulation",
    "⏱️ Parallel execution across all browsers",
    "🔄 Consistent behavior across platforms"
]
add_content_slide(prs, "Browser Support", browsers_content, with_bullet=True)

# ========== SLIDE 5: Architecture ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide, WHITE, LIGHT_GRAY)

# Header
header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
header_fill = header_shape.fill
header_fill.gradient()
header_fill.gradient_stops[0].color.rgb = DARK_BLUE
header_fill.gradient_stops[1].color.rgb = TECH_NAVY
header_shape.line.color.rgb = ACCENT_GREEN
header_shape.line.width = Pt(2)

# Add accent bar
accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), Inches(7.5))
accent_bar.fill.solid()
accent_bar.fill.fore_color.rgb = ACCENT_GREEN
accent_bar.line.color.rgb = ACCENT_GREEN

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
title_p = title_frame.paragraphs[0]
title_p.text = "Project Architecture"
title_p.font.size = Pt(40)
title_p.font.bold = True
title_p.font.color.rgb = WHITE

# Architecture components
arch_items = [
    ("API Layer", "Functions & types for API data\ncinemas, movies, booking, users"),
    ("Fixtures", "Custom Playwright page fixtures\nPre-configured test states"),
    ("Pages", "Page Object Model classes\nBasePage, Forms, Components"),
    ("Tests", "Feature-based test organization\nAuth, booking, account, e2e"),
    ("Reporting", "Allure integration\nVisual test reports")
]

y_pos = 1.5
for i, (title, desc) in enumerate(arch_items):
    # Draw box
    box_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                       Inches(0.5), Inches(y_pos), Inches(9), Inches(0.9))
    box_fill = box_shape.fill
    box_fill.gradient()
    box_fill.gradient_stops[0].color.rgb = LIGHT_BLUE
    box_fill.gradient_stops[1].color.rgb = RGBColor(100, 170, 230)
    box_shape.line.color.rgb = DARK_BLUE
    box_shape.line.width = Pt(2)
    
    # Title in box
    text_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos + 0.05), Inches(2), Inches(0.3))
    tf = text_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    
    # Description in box
    desc_box = slide.shapes.add_textbox(Inches(3), Inches(y_pos + 0.1), Inches(6.2), Inches(0.7))
    dtf = desc_box.text_frame
    dtf.word_wrap = True
    dtf.text = desc
    dtf.paragraphs[0].font.size = Pt(11)
    dtf.paragraphs[0].font.color.rgb = WHITE
    
    y_pos += 1.0

# ========== SLIDE 6: Key Features ==========
left_features = [
    "✅ Page Object Model",
    "✅ Fixture Presets",
    "✅ Parallel Execution",
    "✅ Data Generators"
]
right_features = [
    "✅ Environment Config",
    "✅ Allure Reporting",
    "✅ Screenshot Capture",
    "✅ Multiple Browsers"
]
add_two_column_slide(prs, "Key Features", left_features, right_features)

# ========== SLIDE 7: Test Organization ==========
test_structure = [
    "📁 auth/ - Authentication & login tests",
    "📁 booking/ - Ticket booking workflows",
    "📁 cinemas/ - Cinema information & filtering",
    "📁 account/ - User account management",
    "📁 e2e/ - End-to-end user journeys",
    "📁 responsive/ - Mobile & responsive tests"
]
add_content_slide(prs, "Test Organization", test_structure, with_bullet=True)

# ========== SLIDE 8: Available Test Scripts ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide, WHITE, LIGHT_GRAY)

header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
header_fill = header_shape.fill
header_fill.gradient()
header_fill.gradient_stops[0].color.rgb = DARK_BLUE
header_fill.gradient_stops[1].color.rgb = TECH_NAVY
header_shape.line.color.rgb = ACCENT_GREEN
header_shape.line.width = Pt(2)

# Add accent bar
accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), Inches(7.5))
accent_bar.fill.solid()
accent_bar.fill.fore_color.rgb = ACCENT_GREEN
accent_bar.line.color.rgb = ACCENT_GREEN

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
title_p = title_frame.paragraphs[0]
title_p.text = "Test Execution Commands"
title_p.font.size = Pt(40)
title_p.font.bold = True
title_p.font.color.rgb = WHITE

commands = [
    ("npm run test", "Run all tests in headless mode"),
    ("npm run test:ui", "Interactive UI mode"),
    ("npm run test:headed", "Run with visible browser"),
    ("npm run test:debug", "Debug mode with inspector"),
    ("npm run test:chrome", "Run on Chromium only"),
    ("npm run test:smoke", "Run smoke tests only (@smoke tag)"),
    ("npm run test:regression", "Run regression tests (@regression tag)"),
    ("npm run report", "View Allure test report")
]

y_pos = 1.5
for cmd, desc in commands:
    # Command box
    cmd_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(2.5), Inches(0.35))
    ctf = cmd_box.text_frame
    ctf.text = cmd
    ctf.paragraphs[0].font.size = Pt(12)
    ctf.paragraphs[0].font.bold = True
    ctf.paragraphs[0].font.color.rgb = ACCENT_GREEN
    ctf.paragraphs[0].font.name = "Courier New"
    
    # Description
    desc_box = slide.shapes.add_textbox(Inches(3.5), Inches(y_pos), Inches(5.7), Inches(0.35))
    dtf = desc_box.text_frame
    dtf.word_wrap = True
    dtf.text = desc
    dtf.paragraphs[0].font.size = Pt(12)
    dtf.paragraphs[0].font.color.rgb = DARK_GRAY
    
    y_pos += 0.45

# ========== SLIDE 9: Page Object Model ==========
pom_content = [
    "📄 BasePage - Reusable page actions & locators",
    "📄 BaseForm - Common form operations",
    "📄 LoginPage, RegisterPage, AccountPage, etc.",
    "📄 Components - Reusable UI elements",
    "🎯 Benefits: Maintainability, Reusability, Scalability"
]
add_content_slide(prs, "Page Object Model (POM)", pom_content, with_bullet=True)

# ========== SLIDE 10: Fixtures & Test States ==========
left_fixtures = [
    "🔧 homePage fixture",
    "🔧 loginPage fixture",
    "🔧 registerPage fixture",
    "🔧 accountPage fixture"
]
right_fixtures = [
    "🔑 loggedInHomepage preset",
    "🔑 loggedInUser preset",
    "🔑 Pre-configured states",
    "🔑 Reduced setup code"
]
add_two_column_slide(prs, "Fixtures & Pre-configured States", left_fixtures, right_fixtures)

# ========== SLIDE 11: Best Practice - Fixtures Details ==========
fixtures_detail = [
    "✅ Use standard page fixtures for consistent initialization",
    "✅ Leverage loggedInHomepage for authenticated tests",
    "✅ Use apiClient fixture for direct API testing",
    "✅ Reduce boilerplate code with pre-configured states",
    "✅ Keep fixtures atomic and focused on single responsibility"
]
add_content_slide(prs, "BP: Using Fixtures Effectively", fixtures_detail, with_bullet=True)

# ========== SLIDE 12: Best Practice - Page Objects ==========
pom_detail = [
    "✅ Extend BasePage or BaseForm for consistency",
    "✅ Group locators together at the top of class",
    "✅ Use getBy methods for readable element selection",
    "✅ Create methods for user actions, not just getters",
    "✅ Use appropriate timeouts (SHORT, DEFAULT, LONG)",
    "✅ Avoid duplicating selectors across page objects"
]
add_content_slide(prs, "BP: Page Object Model Design", pom_detail, with_bullet=True)

# ========== SLIDE 13: Best Practice - Testing API ==========
api_bp = [
    "🌐 Use ApiClient for consistent HTTP requests",
    "🌐 Set authorization tokens in beforeEach hooks",
    "🌐 Pass custom headers and timeout overrides",
    "🌐 Validate response status and data in same test",
    "🌐 Test API separately, then with UI for E2E flows"
]
add_content_slide(prs, "BP: API Testing Strategy", api_bp, with_bullet=True)

# ========== SLIDE 14: Best Practice - Writing Tests ==========
test_writing = [
    "📝 Use descriptive test names that explain behavior",
    "📝 Organize tests with test.step() for clarity",
    "📝 Use soft assertions to catch multiple failures",
    "📝 Avoid hard waits - use proper element waits",
    "📝 Extract helper functions to reduce duplication",
    "📝 Tag tests (@smoke, @regression) for selective runs"
]
add_content_slide(prs, "BP: Writing Effective Tests", test_writing, with_bullet=True)

# ========== SLIDE 15: Best Practice - Performance Tips ==========
perf_tips = [
    "⚡ Use efficient selectors: ID, text, or data-testid",
    "⚡ Avoid complex nested CSS selectors",
    "⚡ Reuse fixtures to avoid repeated setup",
    "⚡ Keep tests independent for parallel execution",
    "⚡ Use appropriate timeouts (don't use arbitrary delays)",
    "⚡ Optimize waits for loading spinners and async operations"
]
add_content_slide(prs, "BP: Performance Optimization", perf_tips, with_bullet=True)

# ========== SLIDE 16: Best Practice - Debugging ==========
debug_bp = [
    "🐛 Use npm run test:debug for interactive debugging",
    "🐛 Review logs and screenshots in Allure reports",
    "🐛 Use page.pause() to pause test execution",
    "🐛 Capture screenshots on failures automatically",
    "🐛 Use Playwright Inspector for DOM inspection",
    "🐛 Enable trace viewer for detailed test timeline"
]
add_content_slide(prs, "BP: Debugging Techniques", debug_bp, with_bullet=True)

# ========== SLIDE 17: Best Practice - Common Patterns ==========
patterns = [
    "🔄 Pattern 1: Wait for loading spinner to appear & disappear",
    "🔄 Pattern 2: Retry operations for dynamic content",
    "🔄 Pattern 3: Form submission with validation",
    "🔄 Pattern 4: Extract test data generators",
    "🔄 Pattern 5: Use constants for timeouts & selectors"
]
add_content_slide(prs, "BP: Common Testing Patterns", patterns, with_bullet=True)

# ========== SLIDE 18: Deployment & CI/CD ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide, WHITE, LIGHT_GRAY)

header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1))
header_fill = header_shape.fill
header_fill.gradient()
header_fill.gradient_stops[0].color.rgb = DARK_BLUE
header_fill.gradient_stops[1].color.rgb = TECH_NAVY
header_shape.line.color.rgb = ACCENT_GREEN
header_shape.line.width = Pt(2)

# Add accent bar
accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), Inches(7.5))
accent_bar.fill.solid()
accent_bar.fill.fore_color.rgb = ACCENT_GREEN
accent_bar.line.color.rgb = ACCENT_GREEN

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
title_p = title_frame.paragraphs[0]
title_p.text = "Deployment & Scale"
title_p.font.size = Pt(40)
title_p.font.bold = True
title_p.font.color.rgb = WHITE

features_dp = [
    "🚀 Parallel test execution for faster feedback",
    "⚙️ merge.config.ts for CI/CD optimization",
    "🔄 Sharding strategy for distributed testing",
    "📊 Automatic report generation",
    "🔗 Integration ready for GitHub Actions, Jenkins, etc.",
    "🌍 Cross-environment testing support"
]

y_pos = 1.5
for feature in features_dp:
    box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(8.4), Inches(0.6))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = feature
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = DARK_GRAY
    y_pos += 0.7

# ========== SLIDE 19: Getting Started ==========
getting_started = [
    "1️⃣  Clone the repository",
    "2️⃣  Install dependencies: npm install",
    "3️⃣  Configure environment: .env.qa",
    "4️⃣  Run tests: npm run test",
    "5️⃣  View report: npm run report",
    "📚 Check README.md for detailed documentation"
]
add_content_slide(prs, "Getting Started", getting_started, with_bullet=True)

# ========== SLIDE 20: Thank You ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide, TECH_NAVY, DARK_BLUE)

# Add decorative circles
circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(-0.5), Inches(3.5), Inches(3.5))
circle1.fill.solid()
circle1.fill.fore_color.rgb = ACCENT_GREEN
circle1.fill.transparency = 0.80
circle1.line.color.rgb = ACCENT_GREEN

circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.5), Inches(5), Inches(3), Inches(3))
circle2.fill.solid()
circle2.fill.fore_color.rgb = TECH_CYAN
circle2.fill.transparency = 0.80
circle2.line.color.rgb = TECH_CYAN

# Main thank you message
thank_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
thank_frame = thank_box.text_frame
thank_frame.word_wrap = True
thank_p = thank_frame.paragraphs[0]
thank_p.text = "Thank You!"
thank_p.font.size = Pt(72)
thank_p.font.bold = True
thank_p.font.color.rgb = ACCENT_GREEN
thank_p.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1.2))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.word_wrap = True
subtitle_p = subtitle_frame.paragraphs[0]
subtitle_p.text = "Questions?"
subtitle_p.font.size = Pt(44)
subtitle_p.font.color.rgb = WHITE
subtitle_p.alignment = PP_ALIGN.CENTER

# Footer
footer_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.8))
footer_frame = footer_box.text_frame
footer_p = footer_frame.paragraphs[0]
footer_p.text = "Playwright Framework | End-to-End Testing Solution"
footer_p.font.size = Pt(18)
footer_p.font.color.rgb = TECH_CYAN
footer_p.alignment = PP_ALIGN.CENTER

# Save presentation
output_path = "Playwright_Framework_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation created successfully: {output_path}")
